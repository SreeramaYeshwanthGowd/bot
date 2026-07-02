from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "dpsbot_runtime_flow.pptx"

NAVY = RGBColor(31, 53, 94)
BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(219, 234, 254)
GREEN = RGBColor(22, 163, 74)
LIGHT_GREEN = RGBColor(220, 252, 231)
ORANGE = RGBColor(234, 88, 12)
LIGHT_ORANGE = RGBColor(255, 237, 213)
GRAY = RGBColor(71, 85, 105)
LIGHT_GRAY = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text, size=15, bold=False, color=GRAY, align=PP_ALIGN.CENTER):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.22), Inches(12.6), Inches(0.45))
    set_text(title_box, title, size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.38), Inches(0.68), Inches(12.4), Inches(0.35))
        set_text(subtitle_box, subtitle, size=11, color=GRAY, align=PP_ALIGN.LEFT)


def add_box(slide, x, y, w, h, title, detail, fill, border, title_color=NAVY):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.2)
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.12)
    text_frame.margin_right = Inches(0.12)
    text_frame.margin_top = Inches(0.08)
    text_frame.margin_bottom = Inches(0.08)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = title
    run.font.name = "Aptos"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = title_color

    paragraph = text_frame.add_paragraph()
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = detail
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    return box


def add_arrow(slide, start_x, start_y, end_x, end_y, color=GRAY):
    arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    arrow.line.color.rgb = color
    arrow.line.width = Pt(1.5)
    arrow.line.end_arrowhead = True
    return arrow


def add_step(slide, number, x, y, text, color=BLUE):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.32), Inches(0.32))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = color
    set_text(circle, str(number), size=10, bold=True, color=WHITE)
    label = slide.shapes.add_textbox(x + Inches(0.38), y - Inches(0.02), Inches(4.25), Inches(0.38))
    set_text(label, text, size=10, color=GRAY, align=PP_ALIGN.LEFT)


def create_deck():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    return presentation


def add_overview_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(
        slide,
        "DPSBot Runtime Flow",
        "What happens after a user mentions @DPSBot, from Teams/Web Chat to the final Adaptive Card reply.",
    )

    y = Inches(1.35)
    w = Inches(2.05)
    h = Inches(1.0)
    gap = Inches(0.28)
    x_positions = [Inches(0.35) + i * (w + gap) for i in range(6)]
    boxes = [
        add_box(slide, x_positions[0], y, w, h, "User", "Mentions @DPSBot\nor sends test message", LIGHT_GRAY, GRAY),
        add_box(slide, x_positions[1], y, w, h, "Teams / Web Chat", "Creates Bot Framework\nmessage activity", LIGHT_BLUE, BLUE),
        add_box(slide, x_positions[2], y, w, h, "Azure Bot Service", "Forwards activity to\nconfigured endpoint", LIGHT_BLUE, BLUE),
        add_box(slide, x_positions[3], y, w, h, "function_app.py", "HTTP trigger\n/api/messages", LIGHT_GREEN, GREEN),
        add_box(slide, x_positions[4], y, w, h, "src/dps_bot.py", "on_message_activity\nbuilds reply", LIGHT_GREEN, GREEN),
        add_box(slide, x_positions[5], y, w, h, "src/cards.py", "build_request_card\nreturns form JSON", LIGHT_ORANGE, ORANGE),
    ]

    for i in range(len(boxes) - 1):
        add_arrow(
            slide,
            boxes[i].left + boxes[i].width,
            boxes[i].top + boxes[i].height / 2,
            boxes[i + 1].left,
            boxes[i + 1].top + boxes[i + 1].height / 2,
            BLUE,
        )

    add_arrow(slide, boxes[5].left + boxes[5].width / 2, boxes[5].top + boxes[5].height, boxes[4].left + boxes[4].width / 2, Inches(3.25), ORANGE)
    add_arrow(slide, boxes[4].left + boxes[4].width / 2, Inches(3.25), boxes[3].left + boxes[3].width / 2, Inches(3.25), GREEN)
    add_arrow(slide, boxes[3].left + boxes[3].width / 2, Inches(3.25), boxes[2].left + boxes[2].width / 2, Inches(3.25), GREEN)
    add_arrow(slide, boxes[2].left + boxes[2].width / 2, Inches(3.25), boxes[1].left + boxes[1].width / 2, Inches(3.25), BLUE)

    reply = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.25), Inches(3.65), Inches(6.9), Inches(0.82))
    reply.fill.solid()
    reply.fill.fore_color.rgb = WHITE
    reply.line.color.rgb = GREEN
    set_text(reply, "Reply cycle: card JSON -> Adaptive Card attachment -> Bot Framework response -> Teams/Web Chat renders card", size=13, bold=True, color=NAVY)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(5.15), Inches(12.2), Inches(1.1))
    set_text(
        note,
        "Key config in play: Azure Bot Service endpoint points to the Function URL, and function_app.py uses MicrosoftAppId, MicrosoftAppPassword, and MicrosoftAppTenantId for Bot Framework authentication.",
        size=14,
        color=GRAY,
        align=PP_ALIGN.LEFT,
    )


def add_code_chain_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Code-Level Trigger Chain", "The exact file and method handoff after the message reaches the Function App.")
    steps = [
        ("Azure Function route", "function_app.py\n@app.route(route='messages', methods=['POST'])"),
        ("Read incoming body", "messages(req)\nreq.get_json() or json.loads(raw_body)"),
        ("Deserialize activity", "Activity().deserialize(body)\nauth_header = req.headers['Authorization']"),
        ("Authenticate/process turn", "ADAPTER.process_activity(activity, auth_header, BOT.on_turn)"),
        ("Dispatch to handler", "src/dps_bot.py\nDPSBot.on_message_activity(turn_context)"),
        ("Build card", "src/cards.py\nbuild_request_card()"),
        ("Create reply activity", "MessageFactory.attachment(CardFactory.adaptive_card(card))\nactivity.text + activity.summary"),
        ("Send reply", "turn_context.send_activity(activity)\nBot Service returns it to Teams/Web Chat"),
    ]

    left_x = Inches(0.7)
    right_x = Inches(6.9)
    start_y = Inches(1.25)
    for idx, (title, detail) in enumerate(steps, start=1):
        x = left_x if idx <= 4 else right_x
        y = start_y + (idx - 1 if idx <= 4 else idx - 5) * Inches(1.2)
        box = add_box(slide, x, y, Inches(5.45), Inches(0.86), f"{idx}. {title}", detail, WHITE, BLUE if idx <= 4 else GREEN)
        if idx not in (4, 8):
            add_arrow(slide, box.left + box.width / 2, box.top + box.height, box.left + box.width / 2, box.top + box.height + Inches(0.28), BLUE if idx <= 4 else GREEN)
    add_arrow(slide, Inches(6.15), Inches(4.35), Inches(6.9), Inches(1.68), ORANGE)


def add_plain_english_slide(presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_title(slide, "Final Cycle In Plain English", "Same flow as a quick numbered story for explaining to someone else.")
    text_steps = [
        "User types @DPSBot in Teams or sends a message in Web Chat.",
        "Teams/Web Chat sends a Bot Framework activity to Azure Bot Service.",
        "Azure Bot Service posts that activity to the Function endpoint /api/messages.",
        "function_app.py receives the POST request in messages(req).",
        "function_app.py converts JSON into an Activity and asks the adapter to process it.",
        "The adapter authenticates using app ID, password, and tenant ID, then calls BOT.on_turn.",
        "DPSBot in src/dps_bot.py handles the message and asks src/cards.py for form JSON.",
        "src/cards.py returns the Adaptive Card payload.",
        "src/dps_bot.py wraps the card as an Adaptive Card attachment and sends it.",
        "Azure Bot Service returns the reply to Teams/Web Chat, where the card is rendered.",
    ]
    for idx, text in enumerate(text_steps, start=1):
        x = Inches(0.75) if idx <= 5 else Inches(6.85)
        y = Inches(1.2) + ((idx - 1) % 5) * Inches(0.9)
        add_step(slide, idx, x, y, text, BLUE if idx <= 5 else GREEN)


def build_presentation():
    prs = create_deck()
    add_overview_slide(prs)
    add_code_chain_slide(prs)
    add_plain_english_slide(prs)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(path)